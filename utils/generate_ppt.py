from pptx import Presentation


def create_presentation():
    # Create presentation
    prs = Presentation()

    # Define common layout indices
    TITLE_SLIDE_LAYOUT = 0
    TITLE_AND_CONTENT_LAYOUT = 1
    TWO_CONTENT_LAYOUT = 3
    SECTION_HEADER_LAYOUT = 2

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE_LAYOUT])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "PathWise: Edge-AI for Proactive Road Actor Behavior Prediction"
    subtitle.text = "Real-Time Monocular Traffic Hazard Detection\nSystem Architecture & Prototype Overview"

    # --- Slide 2: Project Objectives ---
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    slide.shapes.title.text = "Project Objectives"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "What does PathWise solve?"

    p = tf.add_paragraph()
    p.text = "• Proactive Warning System: Alerts drivers of imminent hazards before they occur."

    p = tf.add_paragraph()
    p.text = "• Unstructured Environments: Specifically aims to work in non-lane-based traffic (like Indian scenarios)."

    p = tf.add_paragraph()
    p.text = "• Edge Deployment: Designed to run efficiently on Edge Compute devices (e.g., NVIDIA Jetson) using standard dashcam feeds."

    p = tf.add_paragraph()
    p.text = "• Single Camera Solution: Operates purely on monocular (single lens) vision to minimize hardware costs."

    # --- Slide 3: Technical Stack Overview ---
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    slide.shapes.title.text = "Technical Stack"
    tf = slide.shapes.placeholders[1].text_frame

    p = tf.add_paragraph()
    p.text = "Deep Learning Framework: PyTorch & TensorRT"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Computer Vision Layer: OpenCV (cv2)"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Detection Engine: YOLOv10 Nano"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Tracking Algorithm: ByteTrack"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Target Hardware: PC Prototype → NVIDIA Jetson (Deployment)"
    p.level = 0

    # --- Slide 4: Real-Time Perception (Deliverable 1) ---
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    slide.shapes.title.text = "Module 1: Real-Time Perception"
    tf = slide.shapes.placeholders[1].text_frame

    p = tf.add_paragraph()
    p.text = "YOLOv10 Nano Integration"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Ingests 720p/1080p frames rapidly to maintain high FPS rates."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Class Filtering"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Focuses strictly on Cars, Motorcycles, Buses, Pedestrians, Bicycles, and Trucks (COCO class IDs filtering)."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "ByteTrack Implementation"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Solves object occlusion by persisting Object IDs across consecutive frames, preventing system confusion in heavy traffic."
    p.level = 1

    # --- Slide 5: Distance & Velocity Estimation (Deliverable 2) ---
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    slide.shapes.title.text = "Module 2: Monocular Distance Estimation"
    tf = slide.shapes.placeholders[1].text_frame

    p = tf.add_paragraph()
    p.text = "Bird's-Eye View (BEV) Transformation"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Uses Inverse Perspective Mapping (IPM) to warp the 2D image into a flat 'Top-Down' plane. Pixel distances on this plane scale linearly with real-world distance."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Distance Calculation"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Measures the physical distance from the Ego-Vehicle (bottom center) to the bounding box footprint of the tracked targets."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Velocity Buffering"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Maintains a 10-frame sliding window of distances for each target to calculate smooth and accurate Relative Velocities."
    p.level = 1

    # --- Slide 6: Predictive Hazard Logic (Deliverable 3) ---
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    slide.shapes.title.text = "Module 3: Predictive Hazard Engine"
    tf = slide.shapes.placeholders[1].text_frame

    p = tf.add_paragraph()
    p.text = "Time-to-Collision (TTC) Mathematics"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• TTC = Distance / Relative Longitudinal Velocity. Analyzes approaching targets actively closing the gap."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Dynamic Alert Thresholds"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• RED (Critical Risk): TTC ≤ 2.5 seconds."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• YELLOW (Warning): 2.5s < TTC ≤ 4.0 seconds."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• GREEN (Safe): TTC > 4.0 seconds or diverging."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Lateral Cut-In Detection"
    p.level = 0
    p = tf.add_paragraph()
    p.text = '• Monitors cross-frame lateral velocity. Rapid horizontal movement triggers an imminent "CUT-IN" flashing alert.'
    p.level = 1

    # --- Slide 7: Dashboard HUD ---
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    slide.shapes.title.text = "Dashboard Overlay (Heads-Up Display)"
    tf = slide.shapes.placeholders[1].text_frame

    p = tf.add_paragraph()
    p.text = "The system paints a rich, low-latency UI over the dashcam feed:"
    p.level = 0

    p = tf.add_paragraph()
    p.text = (
        "• Hazard-Bounded Boxes: Outlines vehicles in the corresponding Risk Color."
    )
    p.level = 1

    p = tf.add_paragraph()
    p.text = (
        "• Data Tags: Suspends Object ID, Speed (km/h), and TTC metrics above vehicles."
    )
    p.level = 1

    p = tf.add_paragraph()
    p.text = "• Real-Time Statistics: Top status bar reading live FPS and active track counts."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "• BEV Radar Minimap: Plots actors on a spatial 2D grid in the bottom corner, providing a radar-like situational perspective."
    p.level = 1

    # --- Slide 8: Data Logging & Edge Prep (Deliverable 4) ---
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    slide.shapes.title.text = "Module 4: Benchmarking & Optimization"
    tf = slide.shapes.placeholders[1].text_frame

    p = tf.add_paragraph()
    p.text = "System Data Logging"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Autonomously records a high-resolution .csv file asynchronously at 30-frame intervals."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Logs Timestamps, IDs, Distances, Speeds, and specific Hazard Flag states for system accuracy validation."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Edge Deployment Optimization"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Implemented TensorRT export scripts (utils/export_tensorrt.py) to compile the YOLO weights into a highly optimized .engine format with robust FP16 quantization for NVIDIA Jetson architectures."
    p.level = 1

    # Save presentation
    filepath = "C:/Users/anant/Downloads/PathWise_Presentation.pptx"
    prs.save(filepath)
    print(f"Presentation generated at: {filepath}")


if __name__ == "__main__":
    create_presentation()
